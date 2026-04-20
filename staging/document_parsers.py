"""
Document parsers for PDF, PPT, Word, and CSV files.
"""

from typing import Dict
import PyPDF2
from pptx import Presentation
from docx import Document
import csv
import os

class DocumentParser:
    """Multi-format document parser. Supports: PDF, PowerPoint, Word, CSV"""
    
    def parse_file(self, file_path: str) -> Dict[str, str]:
        """Parse document and extract text."""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext in ['.ppt', '.pptx']:
            return self.parse_ppt(file_path)
        elif ext in ['.doc', '.docx']:
            return self.parse_word(file_path)
        elif ext == '.txt':
            return self.parse_text(file_path)
        elif ext == '.csv':
            return self.parse_csv(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    def parse_pdf(self, file_path: str) -> Dict[str, str]:
        """Extract text from PDF."""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            
            return {
                'text': text,
                'metadata': {'pages': len(pdf_reader.pages)}
            }
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
            return {'text': '', 'metadata': {}}
    
    def parse_ppt(self, file_path: str) -> Dict[str, str]:
        """Extract text from PowerPoint."""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return {
                'text': text,
                'metadata': {'slides': len(prs.slides)}
            }
        except Exception as e:
            print(f"Error parsing PPT {file_path}: {e}")
            return {'text': '', 'metadata': {}}
    
    def parse_word(self, file_path: str) -> Dict[str, str]:
        """Extract text from Word document."""
        try:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            
            return {
                'text': text,
                'metadata': {'paragraphs': len(doc.paragraphs)}
            }
        except Exception as e:
            print(f"Error parsing Word {file_path}: {e}")
            return {'text': '', 'metadata': {}}
    
    def parse_text(self, file_path: str) -> Dict[str, str]:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            return {
                'text': text,
                'metadata': {'lines': len(text.split('\n'))}
            }
        except Exception as e:
            print(f"Error parsing text {file_path}: {e}")
            return {'text': '', 'metadata': {}}
    
    def parse_csv(self, file_path: str) -> Dict[str, str]:
        """Extract text from CSV file."""
        try:
            rows = []
            row_count = 0
            
            with open(file_path, 'r', encoding='utf-8') as file:
                csv_reader = csv.DictReader(file)
                headers = csv_reader.fieldnames if csv_reader.fieldnames else []
                
                for row in csv_reader:
                    row_count += 1
                    row_text = " | ".join([f"{k}: {v}" for k, v in row.items() if v])
                    rows.append(row_text)
            
            text = "\n".join(rows)
            
            return {
                'text': text,
                'metadata': {'rows': row_count, 'columns': len(headers)}
            }
        except Exception as e:
            print(f"Error parsing CSV {file_path}: {e}")
            return {'text': '', 'metadata': {}}