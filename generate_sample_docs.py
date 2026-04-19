#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate 100 realistic sample documents (PDF, PPT, Word) for demo.
Creates mock client data across 10 companies.
"""

import os
import sys
import random
from datetime import datetime, timedelta

# Fix Windows console encoding
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt

# Mock client data
CLIENTS = [
    "Flipkart", "HDFC Bank", "Airtel", "Apollo Hospitals", "Tata Motors",
    "Reliance Industries", "Infosys", "Wipro", "ICICI Bank", "Mahindra Group"
]

SERVICE_LINES = ["Strategy", "Operations", "Technology", "Finance", "HR"]
DOCUMENT_TYPES = ["Proposal", "Report", "Analysis", "Roadmap", "Review"]
SENSITIVITIES = ["Internal", "Confidential", "Restricted"]

# Content templates
CONTENT_TEMPLATES = {
    "Strategy": [
        "Strategic expansion proposal for FY{year}. Executive Summary: {client} plans to expand operations across Q1-Q4 {year}.",
        "Market analysis shows strong growth potential in tier-2 and tier-3 cities. Projected revenue increase of 25-30% YoY.",
        "Competitive landscape analysis reveals opportunities in digital transformation and customer experience enhancement.",
        "Five-year strategic roadmap includes technology modernization, talent acquisition, and market expansion initiatives.",
        "Key performance indicators (KPIs) include customer acquisition cost, lifetime value, and market share growth."
    ],
    "Operations": [
        "Operations efficiency report for {client}. Current operational metrics show 98.5% uptime across all facilities.",
        "Process optimization initiatives have reduced operational costs by 15% while maintaining service quality standards.",
        "Supply chain analysis reveals opportunities for automation and predictive maintenance implementation.",
        "Quality assurance metrics indicate 99.2% customer satisfaction with current service delivery standards.",
        "Resource allocation optimization has improved productivity by 22% across all operational units."
    ],
    "Technology": [
        "Technology roadmap for {client}. Cloud migration strategy includes AWS/Azure hybrid deployment by Q4 {year}.",
        "Digital transformation initiatives focus on AI/ML integration, data analytics, and customer-facing applications.",
        "Cybersecurity assessment reveals need for enhanced threat detection and incident response capabilities.",
        "Infrastructure modernization plan includes microservices architecture and containerization strategies.",
        "Technology stack evaluation recommends adoption of modern frameworks and DevOps practices."
    ],
    "Finance": [
        "Financial analysis for {client}. Revenue projections for FY{year} show 18% growth driven by new product launches.",
        "Cost optimization strategies have identified $5M in potential savings through operational efficiencies.",
        "Investment portfolio analysis recommends diversification across emerging markets and technology sectors.",
        "Cash flow projections indicate strong liquidity position with healthy working capital ratios.",
        "Risk assessment highlights currency fluctuations and regulatory changes as key financial considerations."
    ],
    "HR": [
        "Human resources strategy for {client}. Talent acquisition plan targets 500+ hires across technology and operations.",
        "Employee engagement survey results show 85% satisfaction with current workplace culture and benefits.",
        "Learning and development programs have trained 2,000+ employees in digital skills and leadership.",
        "Diversity and inclusion initiatives aim for 40% representation across all management levels by {year}.",
        "Compensation and benefits review recommends competitive adjustments to retain top talent."
    ]
}

# PII data for testing redaction
PII_TEMPLATES = [
    "Contact: john.doe@{domain}.com, +91-9876543210",
    "Project Lead: jane.smith@{domain}.com, +91-8765432109",
    "Account Manager: robert.jones@{domain}.com, +91-7654321098",
    "SSN: 123-45-6789, Credit Card: 4532-1234-5678-9010",
    "Employee ID: EMP{num}, Phone: +91-98765{num}"
]

def generate_pdf(filename, client, service_line, doc_type, content_lines):
    """Generate a PDF document."""
    doc = SimpleDocTemplate(filename, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor='darkblue',
        spaceAfter=30
    )
    title = f"{client} - {service_line} {doc_type}"
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Metadata
    meta_style = styles['Normal']
    story.append(Paragraph(f"<b>Client:</b> {client}", meta_style))
    story.append(Paragraph(f"<b>Service Line:</b> {service_line}", meta_style))
    story.append(Paragraph(f"<b>Document Type:</b> {doc_type}", meta_style))
    story.append(Paragraph(f"<b>Date:</b> {datetime.now().strftime('%Y-%m-%d')}", meta_style))
    story.append(Spacer(1, 0.3*inch))
    
    # Content
    for line in content_lines:
        story.append(Paragraph(line, styles['BodyText']))
        story.append(Spacer(1, 0.1*inch))
    
    # Add PII for testing
    story.append(Spacer(1, 0.3*inch))
    pii = random.choice(PII_TEMPLATES).format(
        domain=client.lower().replace(' ', ''),
        num=random.randint(10000, 99999)
    )
    story.append(Paragraph(f"<b>Contact Information:</b> {pii}", meta_style))
    
    doc.build(story)
    print(f"✓ Created PDF: {filename}")

def generate_ppt(filename, client, service_line, doc_type, content_lines):
    """Generate a PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"{client}"
    subtitle.text = f"{service_line} {doc_type}\n{datetime.now().strftime('%B %Y')}"
    
    # Content slides
    for i, line in enumerate(content_lines[:5], 1):  # Max 5 slides
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"Section {i}"
        text_frame = body_shape.text_frame
        text_frame.text = line
        
        # Add PII on last slide
        if i == len(content_lines[:5]):
            p = text_frame.add_paragraph()
            p.text = f"\n\nContact: {random.choice(PII_TEMPLATES).format(domain=client.lower().replace(' ', ''), num=random.randint(10000, 99999))}"
            p.level = 1
    
    prs.save(filename)
    print(f"✓ Created PPT: {filename}")

def generate_word(filename, client, service_line, doc_type, content_lines):
    """Generate a Word document."""
    doc = Document()
    
    # Title
    title = doc.add_heading(f"{client} - {service_line} {doc_type}", 0)
    
    # Metadata
    doc.add_paragraph(f"Client: {client}")
    doc.add_paragraph(f"Service Line: {service_line}")
    doc.add_paragraph(f"Document Type: {doc_type}")
    doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph()
    
    # Content
    for line in content_lines:
        doc.add_paragraph(line)
    
    # Add PII
    doc.add_paragraph()
    pii = random.choice(PII_TEMPLATES).format(
        domain=client.lower().replace(' ', ''),
        num=random.randint(10000, 99999)
    )
    doc.add_paragraph(f"Contact Information: {pii}")
    
    doc.save(filename)
    print(f"✓ Created Word: {filename}")

def main():
    """Generate 100 sample documents."""
    output_dir = "sample_documents"
    os.makedirs(output_dir, exist_ok=True)
    
    print("=" * 80)
    print("GENERATING 100 SAMPLE DOCUMENTS")
    print("=" * 80)
    print(f"Output directory: {output_dir}/")
    print()
    
    doc_count = 0
    target_count = 100
    
    # Generate mix of PDF (50), PPT (30), Word (20)
    formats = ['pdf'] * 50 + ['ppt'] * 30 + ['docx'] * 20
    random.shuffle(formats)
    
    for fmt in formats:
        client = random.choice(CLIENTS)
        service_line = random.choice(SERVICE_LINES)
        doc_type = random.choice(DOCUMENT_TYPES)
        sensitivity = random.choice(SENSITIVITIES)
        
        # Generate content
        year = random.randint(2024, 2026)
        content_lines = [
            line.format(client=client, year=year)
            for line in random.sample(CONTENT_TEMPLATES[service_line], 3)
        ]
        
        # Filename
        client_short = client.lower().replace(' ', '_')
        filename = f"{output_dir}/{client_short}_{service_line.lower()}_{doc_type.lower()}_{doc_count:03d}.{fmt}"
        
        # Generate document
        try:
            if fmt == 'pdf':
                generate_pdf(filename, client, service_line, doc_type, content_lines)
            elif fmt == 'ppt':
                generate_ppt(filename, client, service_line, doc_type, content_lines)
            elif fmt == 'docx':
                generate_word(filename, client, service_line, doc_type, content_lines)
            
            doc_count += 1
        except Exception as e:
            print(f"⚠️ Error creating {filename}: {e}")
    
    print()
    print("=" * 80)
    print(f"✓ COMPLETE: Generated {doc_count} documents")
    print("=" * 80)
    print()
    print("Documents by type:")
    print(f"  • PDF files: {len([f for f in os.listdir(output_dir) if f.endswith('.pdf')])}")
    print(f"  • PPT files: {len([f for f in os.listdir(output_dir) if f.endswith('.ppt')])}")
    print(f"  • Word files: {len([f for f in os.listdir(output_dir) if f.endswith('.docx')])}")
    print()
    print("Next steps:")
    print("1. Get Cohere API key: https://cohere.com")
    print("2. Set environment variable: COHERE_API_KEY=your_key_here")
    print("3. Install dependencies: pip install -r requirements.txt")
    print("4. Start server: python main.py")
    print()

if __name__ == "__main__":
    main()
